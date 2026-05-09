"""
Convert poster.html → poster.pptx
Screenshot the page at 1920×1080 then embed as a full-slide image.
"""
import asyncio, pathlib, io
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

HTML = pathlib.Path(__file__).parent / "poster.html"
OUT  = pathlib.Path(__file__).parent / "poster.pptx"

W_PX, H_PX = 1920, 1080          # target screenshot resolution
W_EMU = 12192000                  # 1920 × 6350 EMU/px  (13.333 in widescreen)
H_EMU = 6858000                   # 1080 × 6350 EMU/px  (7.5 in widescreen)


async def screenshot() -> bytes:
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page    = await browser.new_page(viewport={"width": W_PX, "height": H_PX})
        await page.goto(HTML.as_uri(), wait_until="networkidle")
        await page.wait_for_timeout(1500)          # let Google Fonts render
        img = await page.screenshot(full_page=False, type="png")
        await browser.close()
        return img


def build_pptx(img_bytes: bytes) -> None:
    prs = Presentation()
    prs.slide_width  = Emu(W_EMU)
    prs.slide_height = Emu(H_EMU)

    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)

    slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        left=Emu(0), top=Emu(0),
        width=Emu(W_EMU), height=Emu(H_EMU),
    )
    prs.save(OUT)
    print(f"Saved → {OUT}")


async def main():
    print("Screenshotting poster.html …")
    img = await screenshot()
    print(f"Screenshot: {len(img)//1024} KB")
    build_pptx(img)


asyncio.run(main())
